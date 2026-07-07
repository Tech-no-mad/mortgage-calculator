import pptx

def inspect_slide5():
    prs = pptx.Presentation(r'C:\Users\karte\Downloads\ServiceNow_x_Accenture_Virtual_Agent.pptx')
    slide = prs.slides[4] # slide 5
    with open('slide5_inspect.txt', 'w', encoding='utf-8') as f:
        f.write("Slide 5 shapes:\n")
        for shp in slide.shapes:
            f.write(f"Shape name: {shp.name}, Type: {shp.shape_type}, Left: {shp.left.inches:.2f}, Top: {shp.top.inches:.2f}\n")
            if shp.has_text_frame:
                f.write(f"  Text: {repr(shp.text_frame.text)}\n")

inspect_slide5()
