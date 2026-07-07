import pptx

def inspect():
    prs = pptx.Presentation(r'C:\Users\karte\Downloads\ServiceNow_x_Accenture_Virtual_Agent.pptx')
    slide = prs.slides[2] # slide 3
    print("Slide 3 shapes:")
    for shp in slide.shapes:
        print(f"Shape name: {shp.name}, Type: {shp.shape_type}, Left: {shp.left.inches:.2f}, Top: {shp.top.inches:.2f}, Width: {shp.width.inches:.2f}, Height: {shp.height.inches:.2f}")
        if shp.shape_type == 6: # Group
            print("  Group contains shapes:")
            for sub in shp.shapes:
                print(f"    Sub-shape name: {sub.name}, Type: {sub.shape_type}, Left: {sub.left.inches:.2f}, Top: {sub.top.inches:.2f}, Width: {sub.width.inches:.2f}, Height: {sub.height.inches:.2f}")
                if sub.has_text_frame:
                    print(f"      Text: {repr(sub.text_frame.text)}")
        elif shp.has_text_frame:
            print(f"  Text: {repr(shp.text_frame.text)}")

if __name__ == '__main__':
    inspect()
