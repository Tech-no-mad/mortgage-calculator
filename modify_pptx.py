import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def apply_replacements():
    in_file = r'C:\Users\karte\Downloads\Virtual_Agent_Self_Service_Incident_Management_Deck.pptx'
    out_file = r'C:\Mortgage Calculator\ServiceNow_x_Accenture_Virtual_Agent.pptx'

    prs = pptx.Presentation(in_file)

    replacements = {
        0: {
            'Intelligent Issue Resolution & Automated Incident Logging': 'Powered by ServiceNow Now Assist × Accenture | ServiceNow Zurich Platform',
            'ATCI Summer Internship Project': 'Presented by: Karteek Tadimalla',
        },
        1: {
            'Why a self-service + auto-incident model is needed': 'The Cost of Manual IT Support',
            'Users struggle to find the right fix': 'Average cost per IT help desk ticket resolved manually: $22',
            'Multiple portals slow issue reporting': 'Rise in ticket volumes since 2020: +16%',
            'Manual ticket logging adds effort': 'Employees attempt to self-serve first: 81%',
            'Delayed support impacts productivity': 'Support teams are stretched thin and overwhelmed',
            'Relevant solutions are difficult to identify quickly.\nUsers often switch between help pages, portals, and support channels.\nIssue intake is inconsistent when tickets are logged manually.\nResolution time increases for repetitive and low-complexity issues.': 'Without self-service, every IT issue — no matter how routine — lands on a human agent, driving up cost and wait times.',
            'Faster first-response time for common IT issues.\nHigher employee satisfaction through guided support.\nLower helpdesk effort for repetitive issue categories.\nA standardized, auditable incident management path.': 'Result: long queues, frustrated employees, agent burnout, and untracked incidents.',
            'A ServiceNow Virtual Agent that handles issue understanding, knowledge guidance, resolution checks, and incident creation in a single conversational flow.': 'AI-powered Virtual Agent (Now Assist) on ServiceNow Zurich enforcing standardized incident creation path.',
            'Captures the user problem in natural language and routes the conversation path.': 'Now Assist conversational interface for natural-language IT support queries.',
            'Surfaces relevant troubleshooting content for self-service before escalation.': 'Instant KB article recommendations before incident creation.',
            'Confirms whether the suggestion resolved the issue or not.': "Enforces use of 'Create Incident - Accenture Support' record producer.",
            'Creates an incident automatically if the problem remains unresolved.': "Routes all incidents to 'ACCENTURE IT SUPPORT' group automatically.",
            'Provides the incident number, assignment path, and confirmation updates.': 'Real-time alerts to the support group and employee closure notification.'
        },
        2: {
            'Logical flow from employee request to incident closure': 'Clean three-layer architecture — all native to ServiceNow Zurich',
            'Provide a seamless self-service experience while automatically routing unresolved issues into governed incident handling.': 'All layers reside natively on the Now Platform with zero external integration required.',
            'Employee / User': 'Employee Interface',
            'Virtual Agent': 'Virtual Agent Chat',
            'Knowledge Base': 'Employee Portal',
            'Resolution Decision': 'AI Orchestration',
            'Incident Creation Logic': 'Custom AI Agent',
            'Incident Table': 'KB Lookup',
            'Assignment Group': 'ITSM Automation',
            'Support Team': 'Auto-Assignment',
            'Notification Services': 'Notifications',
            'Conversation begins in the portal-based Virtual Agent interface.\nKnowledge Base provides first-line guidance for common issues.\nDecision logic determines whether self-service succeeded or escalation is needed.\nIncident record, assignment, and notifications are triggered in a controlled workflow.': 'Virtual Agent Chat is the primary access point.\nNow Assist and Custom AI Agent orchestrate intent and KB lookup.\nRecord Producer ensures standardized data capture.\nAssignment and Notifications are fully automated via ServiceNow native flows.'
        },
        3: {
            'Core ServiceNow building blocks in the solution': 'Implementation Components on ServiceNow Zurich',
            'Conversational interface for user requests, prompts, branching, and guided issue capture.': 'Custom conversation flows covering top IT query categories.',
            'Self-service article recommendations that help employees troubleshoot before escalation.': 'Generative AI skills enabled for knowledge recommendation and intent resolution.',
            'Greeting topic, incident assistance topic, and fallback topic for unsupported paths.': 'Now Assist conversational UI and natural-language understanding.',
            'Orchestrates automation steps such as incident creation, field mapping, and notifications.': "Business rule logic enforces 'Create Incident - Accenture Support' record producer.",
            'Stores incident records and supports lifecycle management, tracking, and updates.': 'All AI-routed incidents land here; zero manual triage required.',
            'Optional request entry point when a specific service request pattern is required.': 'Custom catalog item enforced via AI Agent policy to standardize data.',
            'Reusable business logic for lookups, classification, or specialized processing rules.': "Automated group assignment rule targeting 'ACCENTURE IT SUPPORT'.",
            'Sends confirmations, ticket references, and status updates to users and support teams.': 'Event-driven notifications for support team and automated closure notification.'
        },
        4: {
            'Virtual agent conversation → Knowledge self-service → Auto-incident logging → Closure notification': 'End-to-End Process Flow from query to closed ticket',
            'Employee reports a problem in the portal chat—for example login, access, or device issues.': 'Employee starts conversation with Virtual Agent (Now Assist)',
            'The conversation path identifies the issue type and selects the next guided action.': 'AI Agent searches Knowledge Base, presents relevant articles',
            'The agent presents relevant troubleshooting or how-to guidance for self-service resolution.': 'AI Agent prompts employee to raise an incident if unresolved',
            'The user confirms whether the issue is resolved. Unresolved paths continue automatically.': "Record created automatically via 'Create Incident - Accenture Support'",
            'An incident is created, routed to the right group, and confirmation is sent with the ticket number.': 'Incident auto-assigned to ACCENTURE IT SUPPORT; Team and Employee notified.'
        },
        5: {
            'Measured value across employees, support teams, and the organization': 'Measurable gains in speed, consistency, and employee experience',
            'Immediate assistance through a familiar chat interface.\nReduced wait time for frequent IT issues.\nClear next-step guidance and ticket visibility.\nLess confusion from multi-portal navigation.': '24/7 Self-Service Availability.\nEmployees raise incidents and access knowledge at any hour.\nNo dependency on business hours or agent availability.',
            'Reduced manual ticket logging for repetitive cases.\nConsistent data capture and incident creation standards.\nLower call volume for straightforward resolutions.\nImproved allocation of support effort to higher-complexity work.': '100% Standardized Incident Creation.\nZero Manual Triage Required.\nAll AI-routed incidents are auto-assigned to ACCENTURE IT SUPPORT.',
            'Higher self-service adoption and stronger digital workplace experience.\nFaster issue recovery and improved employee productivity.\nBetter process visibility through standardized incident handling.\nA scalable foundation for future automation use cases.': 'Real-Time Support Team Notifications.\nEvent-driven notification on incident creation.\nReduced manual assignment overhead — agents now focus on resolution.'
        },
        6: {
            'Examples of how the virtual agent operates across common IT needs': 'Live Demo Scenario: VPN Troubleshooting',
            'Scenario 1: Login Issue': 'Initiate Chat',
            'A user reports a sign-in problem. The virtual agent suggests a related knowledge article and basic troubleshooting. If the issue remains unresolved, an incident is created automatically for support follow-up.': 'Ananya opens the Employee Portal and starts a chat with the Virtual Agent.',
            'Scenario 2: Laptop Performance Issue': 'AI Assistance',
            'The virtual agent provides guided steps for device cleanup and performance checks. When the user selects “unable to resolve,” the system logs an incident with the issue details.': 'Now Assist understands intent and surfaces the top 3 VPN troubleshooting KB articles.',
            'Scenario 3: Application Access Request': 'Incident Creation',
            'The chat flow explains the access path, required approvals, or support guidance. If the issue cannot be completed through self-service, the request is routed for support action.': "Issue not resolved — Custom AI Agent creates incident via 'Create Incident - Accenture Support'.",
            'Scenario 4: General IT Support': 'Resolution',
            'The virtual agent serves as a single conversational entry point for broad IT support topics, helping users either resolve simple issues or move smoothly into incident handling.': 'Ticket auto-assigned; Engineer resolves VPN issue; Ananya receives automated closure notification.'
        },
        7: {
            'Current implementation considerations and next-step opportunities': 'Future Roadmap',
            'Challenges Faced': 'Phase 2: Expanded Knowledge',
            'Topic design requires careful branching and fallback handling.\nKnowledge relevance depends on available and well-structured articles.\nIntent recognition must be tuned for varied user phrasing.\nIntegration testing is needed across the complete issue-to-incident flow.': 'Broaden knowledge base coverage; enable multi-turn diagnostic conversations for complex IT issues.',
            'Current Limitations': 'Phase 3: HR & Facilities',
            'The solution is dependent on the breadth and quality of KB content.\nComplex, multi-step issues still require human support involvement.\nUnsupported intents can fall into generic or fallback experiences.\nAdvanced classification may need additional configuration or tuning.': 'Extend Virtual Agent beyond IT to HR and Facilities use cases — all on the same ServiceNow platform.',
            'Future Enhancements': 'Phase 4: Agentic AI',
            'AI-assisted intent classification and smarter topic routing.\nImproved semantic search for better article relevance.\nMulti-channel rollout to collaboration tools or mobile access points.\nPredictive recommendations from recurring incident patterns.': "Autonomous Tier 1 incident resolution with no human intervention, leveraging ServiceNow's evolving AI agent capabilities."
        },
        9: {
            'Project Outcome & Demo Results': 'Final Achievements & Conclusion',
            'What was implemented and what the demonstration proved': 'Built with purpose. Delivered with precision. Powered by ServiceNow Zurich x Accenture.',
            'Delivered a self-service IT support assistant that reduces manual effort while improving employee experience through guided resolution and automated incident management.': 'Demonstrated 100% Standardized Incident Creation and Zero Manual Triage Required using Now Assist on ServiceNow Zurich.'
        }
    }

    for i, slide in enumerate(prs.slides):
        if i in replacements:
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                orig_text = shape.text
                if orig_text in replacements[i]:
                    # To preserve formatting, we replace text in paragraphs and runs
                    new_text = replacements[i][orig_text]
                    shape.text = new_text

        # Add footer to all slides
        txBox = slide.shapes.add_textbox(Inches(8), Inches(7.1), Inches(2), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Karteek Tadimalla"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(128, 128, 128)

    prs.save(out_file)

apply_replacements()
