import pptx
from pptx.util import Inches, Pt

def generate_grid():
    prs = pptx.Presentation()
    # Webdings
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 10 rows, 10 columns
    for row in range(10):
        for col in range(10):
            char_code = 32 + row * 10 + col
            if char_code > 126:
                break
            
            left = Inches(0.5 + col * 1.2)
            top = Inches(0.5 + row * 0.6)
            
            tb = slide.shapes.add_textbox(left, top, Inches(1.1), Inches(0.5))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            
            # Show ascii char, code, and then the webdings char
            run1 = p.add_run()
            run1.text = f"{char_code}:{chr(char_code)}="
            run1.font.name = "Arial"
            run1.font.size = Pt(8)
            
            run2 = p.add_run()
            run2.text = f"{chr(char_code)}"
            run2.font.name = "Webdings"
            run2.font.size = Pt(14)
            
    prs.save('webdings_grid.pptx')
    
    # Wingdings
    prs2 = pptx.Presentation()
    slide2 = prs2.slides.add_slide(prs2.slide_layouts[6])
    for row in range(10):
        for col in range(10):
            char_code = 32 + row * 10 + col
            if char_code > 126:
                break
            left = Inches(0.5 + col * 1.2)
            top = Inches(0.5 + row * 0.6)
            tb = slide2.shapes.add_textbox(left, top, Inches(1.1), Inches(0.5))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            
            run1 = p.add_run()
            run1.text = f"{char_code}:{chr(char_code)}="
            run1.font.name = "Arial"
            run1.font.size = Pt(8)
            
            run2 = p.add_run()
            run2.text = f"{chr(char_code)}"
            run2.font.name = "Wingdings"
            run2.font.size = Pt(14)
    prs2.save('wingdings_grid.pptx')

generate_grid()
