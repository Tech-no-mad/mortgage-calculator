import pptx

def inspect_slide10():
    prs = pptx.Presentation(r'C:\Users\karte\Downloads\ServiceNow_x_Accenture_Virtual_Agent.pptx')
    slide = prs.slides[9] # slide 10
    with open('slide10_inspect.txt', 'w', encoding='utf-8') as f:
        f.write("Slide 10 shapes:\n")
        for shp in slide.shapes:
            f.write(f"Shape name: {shp.name}, Type: {shp.shape_type}\n")
            if shp.has_text_frame:
                f.write(f"  Text: {repr(shp.text_frame.text)}\n")

inspect_slide10()
