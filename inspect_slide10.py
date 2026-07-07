import pptx

def inspect_slide10():
    prs = pptx.Presentation(r'C:\Users\karte\Downloads\ServiceNow_x_Accenture_Virtual_Agent.pptx')
    slide = prs.slides[9] # slide 10
    print("Slide 10 shapes:")
    for shp in slide.shapes:
        print(f"Shape name: {shp.name}, Type: {shp.shape_type}")
        if shp.has_text_frame:
            print(f"  Text: {repr(shp.text_frame.text)}")

inspect_slide10()
