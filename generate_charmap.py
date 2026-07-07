import pptx
from pptx.util import Inches, Pt

def generate_charmap():
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(6.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    # Let's write Wingdings characters
    p = tf.paragraphs[0]
    p.text = "Wingdings Font Map (32-126):\n"
    p.font.size = Pt(14)
    
    for i in range(32, 127):
        p = tf.add_paragraph()
        run1 = p.add_run()
        run1.text = f"{chr(i)} ({i}): "
        run1.font.name = "Arial"
        run1.font.size = Pt(12)
        
        run2 = p.add_run()
        run2.text = f" {chr(i)} "
        run2.font.name = "Wingdings"
        run2.font.size = Pt(16)
        
    prs.save('wingdings_charmap.pptx')
    
    # Also create Webdings
    prs2 = pptx.Presentation()
    slide2 = prs2.slides.add_slide(prs2.slide_layouts[6])
    tb2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(6.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Webdings Font Map (32-126):\n"
    p2.font.size = Pt(14)
    
    for i in range(32, 127):
        p = tf2.add_paragraph()
        run1 = p.add_run()
        run1.text = f"{chr(i)} ({i}): "
        run1.font.name = "Arial"
        run1.font.size = Pt(12)
        
        run2 = p.add_run()
        run2.text = f" {chr(i)} "
        run2.font.name = "Webdings"
        run2.font.size = Pt(16)
        
    prs2.save('webdings_charmap.pptx')
    print("Saved charmaps!")

generate_charmap()
