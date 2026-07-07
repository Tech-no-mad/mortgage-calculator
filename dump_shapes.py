import pptx

def dump_shapes(file_path, out_file):
    out_file.write(f'\n--- {file_path} ---\n')
    try:
        prs = pptx.Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            out_file.write(f'Slide {i+1}:\n')
            for j, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text') and shape.text.strip():
                    text = shape.text.replace('\n', '\\n')
                    out_file.write(f'  Shape {j} ({shape.name}): {text}\n')
    except Exception as e:
        out_file.write(f'Error: {e}\n')

with open(r'C:\Mortgage Calculator\ppt_shapes.txt', 'w', encoding='utf-8') as f:
    dump_shapes(r'C:\Users\karte\Downloads\Virtual_Agent_Self_Service_Incident_Management_Deck.pptx', f)
