import pptx
from pptx.util import Inches

def print_titles():
    prs = pptx.Presentation(r'C:\Users\karte\Downloads\ServiceNow_x_Accenture_Virtual_Agent.pptx')
    for i, slide in enumerate(prs.slides):
        title = ""
        for shp in slide.shapes:
            if shp.has_text_frame and shp.text_frame.text:
                # check if it's the title
                if shp.top < Inches(1.5) and shp.width > Inches(5):
                    title = shp.text_frame.text.split('\n')[0]
                    break
        print(f"Slide {i+1}: {repr(title)}")

print_titles()
