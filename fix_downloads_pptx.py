import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re
import copy
import traceback

ACCENTURE_PURPLE = RGBColor(161, 0, 255)
DARK_GRAY = RGBColor(64, 64, 64)
MEDIUM_GRAY = RGBColor(128, 128, 128)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(20, 20, 20)

def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == 6: # Group
            for sub_shape in iter_shapes(shape.shapes):
                yield sub_shape
        else:
            yield shape

def fix_presentation(input_path, output_path):
    prs = pptx.Presentation(input_path)

    # 6. Split Slide 2 (index 1) into Slide 2a (The Problem) and Slide 2b (Proposed Solution)
    slide2 = prs.slides[1]
    slide3 = prs.slides.add_slide(prs.slide_layouts[6]) 
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    xml_slides.remove(slides_list[-1])
    xml_slides.insert(2, slides_list[-1])
    slide3 = prs.slides[2]
    
    # Deepcopy all shapes from slide 2 to slide 3
    for shp in slide2.shapes:
        el = shp.element
        new_el = copy.deepcopy(el)
        slide3.shapes._spTree.append(new_el)
        
    # Clean up Slide 2 (Slide 2a: The Problem)
    to_delete_s2 = []
    for shp in list(slide2.shapes):
        if shp.left > Inches(3.5) and shp.top > Inches(3.4):
            to_delete_s2.append(shp)
    for shp in to_delete_s2:
        try: shp.element.getparent().remove(shp.element)
        except: pass
        
    # Clean up Slide 3 (Slide 2b: Proposed Solution)
    to_delete_s3 = []
    for shp in list(slide3.shapes):
        if shp.left > Inches(3.5) and shp.top < Inches(3.4):
            to_delete_s3.append(shp)
    for shp in to_delete_s3:
        try: shp.element.getparent().remove(shp.element)
        except: pass
        
    # Move bottom content on Slide 3 up and change titles
    for shp in list(slide3.shapes):
        if shp.left > Inches(3.5) and shp.top > Inches(3.4):
            shp.top -= Inches(2.55)
        # Rename sidebar titles on Slide 3
        if shp.left < Inches(3.5) and shp.has_text_frame:
            txt = shp.text_frame.text.strip()
            if txt == "Problem Statement":
                shp.text_frame.text = "Proposed Solution"
                for p in shp.text_frame.paragraphs:
                    p.font.size = Pt(28)
                    p.font.bold = True
                    p.font.color.rgb = ACCENTURE_PURPLE
            elif txt == "The Cost of Manual IT Support":
                shp.text_frame.text = "Standardized Flow"
                for p in shp.text_frame.paragraphs:
                    p.font.size = Pt(18)
                    p.font.color.rgb = DARK_GRAY
            elif "KEY DRIVERS" in txt:
                shp.text_frame.text = "FLOW STEPS"
                for p in shp.text_frame.paragraphs:
                    p.font.size = Pt(12)
                    p.font.bold = True
                    p.font.color.rgb = DARK_GRAY
            elif txt.startswith("Average cost") or txt.startswith("Rise in ticket") or txt.startswith("Employees attempt") or txt.startswith("Support teams"):
                try: shp.element.getparent().remove(shp.element)
                except: pass
            elif txt == "•":
                try: shp.element.getparent().remove(shp.element)
                except: pass
            elif shp.shape_type == 1 and shp.width < Inches(0.5) and shp.top > Inches(3.0):
                try: shp.element.getparent().remove(shp.element)
                except: pass

    # 1. Slide 11 Leftover note & overlap (Closing slide is slide 11, index 10)
    slide11 = prs.slides[10]
    to_delete = []
    for shp in list(iter_shapes(slide11.shapes)):
        if shp.has_text_frame:
            txt = shp.text_frame.text
            if "Reference-inspired" in txt or "Built with purpose" in txt:
                to_delete.append(shp)
    for shp in to_delete:
        try: shp.element.getparent().remove(shp.element)
        except: pass
        
    tb = slide11.shapes.add_textbox(Inches(1.5), Inches(6.5), Inches(7), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Ready to accelerate your IT support transformation."
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENTURE_PURPLE

    # 8. Strip underlines (including line shapes under titles)
    for slide in prs.slides:
        to_delete = []
        for shp in list(iter_shapes(slide.shapes)):
            if shp.has_text_frame:
                for para in shp.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.underline = False
            # Check for thin lines/rectangles under subtitles
            if shp.shape_type == 9: # msoLine
                if shp.top < Inches(2.0):
                    to_delete.append(shp)
            elif shp.shape_type == 1 and getattr(shp, "auto_shape_type", None) == MSO_SHAPE.RECTANGLE:
                if shp.height < Inches(0.15) and shp.top < Inches(2.0): 
                    to_delete.append(shp)
        for shp in to_delete:
            try: shp.element.getparent().remove(shp.element)
            except: pass

    # 2 & 10. Remove bare "Karteek Tadimalla" and purple bottom bar
    for slide in prs.slides:
        to_delete = []
        for shp in list(slide.shapes): 
            if shp.has_text_frame:
                txt = shp.text_frame.text.strip()
                if txt == "Karteek Tadimalla" or " | " in txt:
                    to_delete.append(shp)
            if getattr(shp, 'shape_type', None) == 1: 
                if shp.width > Inches(8) and shp.top > Inches(6.8):
                    to_delete.append(shp)
        for shp in to_delete:
            try: shp.element.getparent().remove(shp.element)
            except: pass

    # 5. Fix Slide 6 (Process Flow, index 5) left sidebar overlaps
    slide6_process = prs.slides[5]
    for shp in list(slide6_process.shapes):
        if shp.left < Inches(3.5):
            if shp.has_text_frame:
                txt = shp.text_frame.text.strip()
                # Check subtitle first to avoid matching "Process Flow" in subtitle
                if "End-to-End Process" in txt:
                    shp.top = Inches(1.10)
                    for p in shp.text_frame.paragraphs:
                        p.font.size = Pt(11)
                elif "Process Flow" in txt:
                    shp.width = Inches(2.8)
                elif "KEY PRINCIPLES" in txt:
                    shp.top = Inches(2.20)
                elif txt == "✓":
                    if abs(shp.top.inches - 2.58) < 0.3: shp.top = Inches(2.70)
                    elif abs(shp.top.inches - 3.38) < 0.3: shp.top = Inches(3.50)
                    elif abs(shp.top.inches - 4.18) < 0.3: shp.top = Inches(4.30)
                    elif abs(shp.top.inches - 4.98) < 0.3: shp.top = Inches(5.10)
                elif "helpdesk call" in txt or "No helpdesk" in txt:
                    shp.top = Inches(2.70)
                elif "ticket creation" in txt or "No manual" in txt:
                    shp.top = Inches(3.50)
                elif "system navigation" in txt or "No system" in txt:
                    shp.top = Inches(4.30)
                elif "technical knowledge" in txt or "No technical" in txt:
                    shp.top = Inches(5.10)
            elif shp.shape_type == 1 and shp.width < Inches(0.5):
                if abs(shp.top.inches - 2.55) < 0.3: shp.top = Inches(2.70)
                elif abs(shp.top.inches - 3.35) < 0.3: shp.top = Inches(3.50)
                elif abs(shp.top.inches - 4.15) < 0.3: shp.top = Inches(4.30)
                elif abs(shp.top.inches - 4.95) < 0.3: shp.top = Inches(5.10)

    # 7. Fix Architecture diagram on Slide 4 (index 3)
    slide4 = prs.slides[3]
    for shp in list(slide4.shapes):
        if shp.has_text_frame:
            txt = shp.text_frame.text.strip()
            if "Architecture Notes" in txt or "Virtual Agent Chat is the primary" in txt:
                shp.top -= Inches(0.4)
            if "Portal conversation" in txt or "→" in txt:
                shp.top = Inches(6.8) 
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.color.rgb = DARK_GRAY

    # 3. Boilerplate captions
    # Slide 7 (Benefits, index 6)
    slide7_benefits = prs.slides[6]
    caps_6 = [shp for shp in list(iter_shapes(slide7_benefits.shapes)) if shp.has_text_frame and "Higher consistency and faster response" in shp.text_frame.text]
    caps_6.sort(key=lambda x: x.left)
    r6 = ["Instant 24/7 Support", "Zero-Touch Routing", "Scalable Automation"]
    for i, shp in enumerate(caps_6):
        if i < len(r6): shp.text_frame.text = r6[i]

    # Slide 8 (Scenarios, index 7)
    slide8_scenarios = prs.slides[7]
    caps_7 = [shp for shp in list(iter_shapes(slide8_scenarios.shapes)) if shp.has_text_frame and "Self-service" in shp.text_frame.text and "Escalation" in shp.text_frame.text]
    # Sort correctly in 2D reading order (Row 1 left-to-right, then Row 2 left-to-right)
    caps_7.sort(key=lambda x: (round(x.top.inches), x.left.inches))
    r7 = ["Guided Assistance", "Knowledge Deflection", "Automated Ticketing", "Seamless Follow-up"]
    for i, shp in enumerate(caps_7):
        if i < len(r7): shp.text_frame.text = r7[i]

    # Slide 9 (Roadmap, index 8)
    slide9_roadmap = prs.slides[8]
    caps_8 = [shp for shp in list(iter_shapes(slide9_roadmap.shapes)) if shp.has_text_frame and "Design for iteration and scale" in shp.text_frame.text]
    caps_8.sort(key=lambda x: x.left)
    r8 = ["Wider Reach", "Cross-Domain Support", "Autonomous IT"]
    for i, shp in enumerate(caps_8):
        if i < len(r8): shp.text_frame.text = r8[i]

    # 4. Text codes to Webdings/Wingdings icons
    codes_map = {
        "USR": ("Webdings", "u"),
        "EMP": ("Webdings", "u"),
        "CS":  ("Webdings", "u"),
        "LM":  ("Webdings", "u"),
        "VA":  ("Webdings", "("),
        "D":   ("Webdings", "("),
        "KB":  ("Wingdings", "&"),
        "CT":  ("Wingdings", "&"),
        "INC": ("Webdings", "t"),
        "FD":  ("Wingdings", "R"),
        "IT":  ("Wingdings", "R"),
        "IN":  ("Wingdings", "R"),
        "SI":  ("Wingdings", "R"),
        "DB":  ("Wingdings", "R"),
        "ORG": ("Webdings", "C"),
        "TP":  ("Wingdings", ":"),
        "NT":  ("Wingdings", ":"),
        "CH":  ("Webdings", "r"),
        "FX":  ("Webdings", "a"),
        "R":   ("Webdings", "a")
    }
    for slide in prs.slides:
        for shp in list(iter_shapes(slide.shapes)):
            if shp.has_text_frame:
                txt = shp.text_frame.text.strip()
                if txt in codes_map:
                    font_name, char = codes_map[txt]
                    shp.text_frame.text = char
                    for p in shp.text_frame.paragraphs:
                        p.alignment = 2 # center
                        for r in p.runs:
                            r.font.name = font_name
                            r.font.size = Pt(18)
                            r.font.color.rgb = WHITE

    # 5. Colors & Dark theme (Slide 1 and 11)
    for i, slide in enumerate(prs.slides):
        slide_idx = i + 1
        
        # Dark Theme for Title (Slide 1) and Closing (Slide 11)
        if slide_idx in [1, 11]:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = BLACK
            
            for shp in list(slide.shapes):
                if shp.shape_type == 1:
                    if shp.width > Inches(3.0) and shp.left < Inches(0.5):
                        shp.fill.solid()
                        shp.fill.fore_color.rgb = ACCENTURE_PURPLE
                        shp.line.fill.background()
                    elif shp.shape_type == 1 and shp.width > Inches(1.0) and shp.left > Inches(3.5):
                        shp.fill.solid()
                        shp.fill.fore_color.rgb = ACCENTURE_PURPLE
                        shp.line.fill.background()
                        
            # Set text on dark slides to white
            for shp in list(iter_shapes(slide.shapes)):
                if shp.has_text_frame:
                    for p in shp.text_frame.paragraphs:
                        for r in p.runs:
                            r.font.color.rgb = WHITE
                            
            # Fix Slide 11 card descriptions truncation: make font size Pt(10.5)
            if slide_idx == 11:
                # Adjust card and description heights
                for shp in list(slide.shapes):
                    if shp.name in ["Rounded Rectangle 26", "Rounded Rectangle 31"]:
                        shp.height = Inches(2.40)
                    elif shp.name in ["TextBox 29", "TextBox 34"]:
                        shp.top = Inches(1.20)
                    elif shp.name in ["TextBox 30", "TextBox 35"]:
                        shp.top = Inches(1.55)
                        shp.height = Inches(1.85)
                
                # Apply 10.5pt font size
                for shp in list(iter_shapes(slide.shapes)):
                    if shp.has_text_frame and shp.left > Inches(3.5) and len(shp.text_frame.text) > 50:
                        for p in shp.text_frame.paragraphs:
                            p.font.size = Pt(10.5)
                            for r in p.runs:
                                r.font.size = Pt(10.5)
                            
        # Style Slide 4 Architecture Flow matching Legend Colors
        elif slide_idx == 4:
            for shp in list(iter_shapes(slide.shapes)):
                if shp.has_text_frame:
                    txt = shp.text_frame.text.strip()
                    is_accenture = txt in ["Employee Interface", "Virtual Agent Chat", "Employee Portal", "Employee Confirmation", "Notifications"]
                    is_medium = txt in ["AI Orchestration", "Auto-Assignment", "ITSM Automation"]
                    is_dark = txt in ["Custom AI Agent", "KB Lookup"]
                    
                    target_color = None
                    if is_accenture: target_color = ACCENTURE_PURPLE
                    elif is_medium: target_color = MEDIUM_GRAY
                    elif is_dark: target_color = DARK_GRAY
                    
                    if target_color:
                        for other in list(slide.shapes):
                            if other.shape_type == 1 and abs(other.left.inches - shp.left.inches) < 0.1:
                                if abs(other.top.inches - shp.top.inches + 0.2) < 0.3:
                                    if other.height < Inches(0.2): 
                                        other.fill.solid()
                                        other.fill.fore_color.rgb = target_color
            
            # Color Legend boxes exactly
            for shp in list(slide.shapes):
                if shp.shape_type == 1 and shp.width < Inches(0.3) and shp.height < Inches(0.3):
                    if abs(shp.left.inches - 8.55) < 0.1:
                        if abs(shp.top.inches - 5.38) < 0.1: 
                            shp.fill.solid()
                            shp.fill.fore_color.rgb = ACCENTURE_PURPLE
                        elif abs(shp.top.inches - 5.68) < 0.1: 
                            shp.fill.solid()
                            shp.fill.fore_color.rgb = MEDIUM_GRAY
                        elif abs(shp.top.inches - 5.98) < 0.1: 
                            shp.fill.solid()
                            shp.fill.fore_color.rgb = DARK_GRAY
                            
        # Slide 10: Testing Table contrast fix
        elif slide_idx == 10:
            # First, change all non-white shapes' fill color to ACCENTURE_PURPLE
            for shp in list(iter_shapes(slide.shapes)):
                if getattr(shp, 'fill', None) and getattr(shp.fill, 'type', None) == 1:
                    try:
                        color_rgb = shp.fill.fore_color.rgb
                        if color_rgb != WHITE and color_rgb != RGBColor(255, 255, 255):
                            shp.fill.fore_color.rgb = ACCENTURE_PURPLE
                    except:
                        pass
                        
            # Second, find all text boxes overlapping a purple shape and make their text white
            for tb in list(iter_shapes(slide.shapes)):
                if tb.shape_type == 17 and tb.has_text_frame:
                    is_on_purple = False
                    for rect in list(iter_shapes(slide.shapes)):
                        if rect.shape_type == 1 and getattr(rect, 'fill', None) and getattr(rect.fill, 'type', None) == 1:
                            try:
                                if rect.fill.fore_color.rgb == ACCENTURE_PURPLE:
                                    if abs(rect.left - tb.left) < Inches(0.2) and abs(rect.top - tb.top) < Inches(0.2):
                                        is_on_purple = True
                                        break
                            except:
                                pass
                    if is_on_purple:
                        for p in tb.text_frame.paragraphs:
                            if p.text.strip():
                                if not p.runs:
                                    # Create a run if it doesn't exist to ensure we can apply font styling
                                    txt = p.text
                                    p.text = ""
                                    run = p.add_run()
                                    run.text = txt
                                for r in p.runs:
                                    r.font.color.rgb = WHITE
                            
        # For all other light slides, change theme purple to ACCENTURE_PURPLE
        else:
            for shp in list(iter_shapes(slide.shapes)):
                if getattr(shp, 'fill', None) and shp.fill.type == 1:
                    try:
                        color_rgb = shp.fill.fore_color.rgb
                        if color_rgb != WHITE and color_rgb != RGBColor(255, 255, 255):
                            shp.fill.fore_color.rgb = ACCENTURE_PURPLE
                    except:
                        pass

    # 9. Renumbering cascade & proper footer
    for i, slide in enumerate(prs.slides):
        slide_idx = i + 1
        
        # Add footer (only on light slides)
        if slide_idx not in [1, 11]:
            tb = slide.shapes.add_textbox(Inches(10.0), Inches(7.0), Inches(3.0), Inches(0.3))
            p = tb.text_frame.paragraphs[0]
            p.text = f"Karteek Tadimalla | Slide {slide_idx}"
            p.font.size = Pt(9)
            p.font.color.rgb = MEDIUM_GRAY
            
        # Update chip badge slide number
        for shp in list(iter_shapes(slide.shapes)):
            if shp.has_text_frame:
                txt = shp.text_frame.text.strip()
                if re.match(r'^0?\d+$', txt) and len(txt) <= 2:
                    if shp.left < Inches(3.0) and shp.top < Inches(1.0):
                        p_badge = shp.text_frame.paragraphs[0]
                        p_badge.text = f"{slide_idx:02d}"
                        for run in p_badge.runs:
                            run.font.color.rgb = ACCENTURE_PURPLE 
                            run.font.bold = True

    prs.save(output_path)
    print("DONE SUCCESSFULLY!")

if __name__ == '__main__':
    try:
        fix_presentation(
            r'C:\Users\karte\Downloads\ServiceNow_x_Accenture_Virtual_Agent.pptx',
            r'C:\Users\karte\Downloads\ServiceNow_x_Accenture_Virtual_Agent.pptx'
        )
    except Exception as e:
        traceback.print_exc()
