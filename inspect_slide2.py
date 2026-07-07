import pptx

def inspect_slide2():
    prs = pptx.Presentation(r'C:\Users\karte\Downloads\ServiceNow_x_Accenture_Virtual_Agent.pptx')
    slide = prs.slides[1] # slide 2
    with open('slide2_inspect.txt', 'w', encoding='utf-8') as f:
        f.write("Slide 2 shapes:\n")
        for shp in slide.shapes:
            f.write(f"Shape name: {shp.name}, Type: {shp.shape_type}, Left: {shp.left.inches:.2f}, Top: {shp.top.inches:.2f}, Width: {shp.width.inches:.2f}, Height: {shp.height.inches:.2f}\n")
            if shp.has_text_frame:
                f.write(f"  Text: {repr(shp.text_frame.text)}\n")

inspect_slide2()
