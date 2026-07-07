import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re

ACCENTURE_PURPLE = RGBColor(161, 0, 255)
DARK_GRAY = RGBColor(64, 64, 64)
WHITE = RGBColor(255, 255, 255)

def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == 6: # msoGroup
            for sub_shape in iter_shapes(shape.shapes):
                yield sub_shape
        else:
            yield shape

def fix_presentation(input_path, output_path):
    prs = pptx.Presentation(input_path)

    # 1. Slide 10
    slide10 = prs.slides[9]
    to_delete = []
    for shp in iter_shapes(slide10.shapes):
        if shp.has_text_frame:
            txt = shp.text_frame.text
            if "Reference-inspired layout" in txt or "Built with purpose" in txt:
                to_delete.append(shp)
    for shp in to_delete:
        try:
            shp.element.getparent().remove(shp.element)
        except:
            pass
        
    tb = slide10.shapes.add_textbox(Inches(1.5), Inches(6.5), Inches(7), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Ready to accelerate your IT support transformation."
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENTURE_PURPLE

    # 8. Strip underlines (including line shapes under titles)
    for slide in prs.slides:
        to_delete = []
        for shp in iter_shapes(slide.shapes):
            if shp.has_text_frame:
                for para in shp.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.underline = False
            if shp.shape_type == 9: # msoLine
                to_delete.append(shp)
            elif shp.shape_type == 1 and getattr(shp, "auto_shape_type", None) == MSO_SHAPE.RECTANGLE:
                if shp.height < Inches(0.1): 
                    to_delete.append(shp)
        for shp in to_delete:
            try: shp.element.getparent().remove(shp.element)
            except: pass

    # 2 & 10. Remove bare "Karteek Tadimalla" and purple bottom bar
    for slide in prs.slides:
        to_delete = []
        for shp in slide.shapes: # Only at top level for bottom bar
            if shp.has_text_frame:
                txt = shp.text_frame.text.strip()
                if txt == "Karteek Tadimalla" or " | " in txt:
                    to_delete.append(shp)
            if getattr(shp, 'shape_type', None) == 1: 
                if shp.width > Inches(8) and shp.top > Inches(6):
                    to_delete.append(shp)
        for shp in to_delete:
            try: shp.element.getparent().remove(shp.element)
            except: pass

    # 6. Split Slide 2
    slide2 = prs.slides[1]
    slide3 = prs.slides.add_slide(prs.slide_layouts[6]) 
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    xml_slides.remove(slides_list[-1])
    xml_slides.insert(2, slides_list[-1])
    slide3 = prs.slides[2]
    
    to_move = []
    for shp in slide2.shapes:
        if shp.top > Inches(4.5): 
            to_move.append(shp)
    for shp in to_move:
        shp.top -= Inches(2.5) 
        slide3.shapes._spTree.append(shp.element)
        
    tb = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Proposed Solution"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENTURE_PURPLE

    # 7. Fix Architecture diagram on Slide 4 (formerly 3)
    slide4 = prs.slides[3]
    for shp in iter_shapes(slide4.shapes):
        if shp.has_text_frame:
            txt = shp.text_frame.text.strip()
            # Top row reset positions
            if txt in ["Employee\nInterface", "Employee Interface"]: shp.left, shp.width = Inches(0.5), Inches(1.2)
            if txt in ["Virtual Agent\nChat", "Virtual Agent Chat"]: shp.left, shp.width = Inches(1.8), Inches(1.2)
            if txt in ["Employee\nPortal", "Employee Portal"]: shp.left, shp.width = Inches(3.1), Inches(1.2)
            if txt in ["AI\nOrchestration", "AI Orchestration"]: shp.left, shp.width = Inches(4.4), Inches(1.2)
            if txt in ["Custom AI\nAgent", "Custom AI Agent"]: shp.left, shp.width = Inches(5.7), Inches(1.2)
            # Bottom row
            if txt in ["Employee\nConfirmation", "Employee Confirmation"]: shp.left, shp.width = Inches(0.5), Inches(1.2)
            if txt == "Notifications": shp.left, shp.width = Inches(1.8), Inches(1.2)
            if txt in ["Auto-\nAssignment", "Auto-Assignment"]: shp.left, shp.width = Inches(3.1), Inches(1.2)
            if txt in ["ITSM\nAutomation", "ITSM Automation"]: shp.left, shp.width = Inches(4.4), Inches(1.2)
            if txt == "KB Lookup": shp.left, shp.width = Inches(5.7), Inches(1.2)
            
            if "Architecture Notes" in txt: 
                shp.top -= Inches(0.5)
                if getattr(shp, 'fill', None):
                    shp.fill.solid()
                    shp.fill.fore_color.rgb = WHITE
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.color.rgb = DARK_GRAY
            
            if "Portal conversation" in txt or "→" in txt:
                shp.top -= Inches(0.5)
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.color.rgb = DARK_GRAY

    # 3. Boilerplate captions
    slide7 = prs.slides[6]
    caps = [shp for shp in iter_shapes(slide7.shapes) if shp.has_text_frame and "Higher consistency and faster response" in shp.text_frame.text]
    caps.sort(key=lambda x: x.left)
    r7 = ["Instant 24/7 Support", "Zero-Touch Routing", "Scalable Automation"]
    for i, shp in enumerate(caps):
        if i < len(r7): shp.text_frame.text = r7[i]

    slide8 = prs.slides[7]
    caps = [shp for shp in iter_shapes(slide8.shapes) if shp.has_text_frame and "Self-service" in shp.text_frame.text and "Escalation" in shp.text_frame.text]
    caps.sort(key=lambda x: x.left)
    r8 = ["Guided Assistance", "Knowledge Deflection", "Automated Ticketing", "Seamless Follow-up"]
    for i, shp in enumerate(caps):
        if i < len(r8): shp.text_frame.text = r8[i]

    slide9 = prs.slides[8]
    caps = [shp for shp in iter_shapes(slide9.shapes) if shp.has_text_frame and "Design for iteration and scale" in shp.text_frame.text]
    caps.sort(key=lambda x: x.left)
    r9 = ["Wider Reach", "Cross-Domain Support", "Autonomous IT"]
    for i, shp in enumerate(caps):
        if i < len(r9): shp.text_frame.text = r9[i]

    # 4. Text codes to shapes
    codes = ["USR", "VA", "KB", "INC", "IT", "EMP", "ORG", "DB"]
    shape_types = [MSO_SHAPE.GEAR_6, MSO_SHAPE.CLOUD, MSO_SHAPE.HEART, MSO_SHAPE.STAR_5_POINT, MSO_SHAPE.DIAMOND, MSO_SHAPE.HEXAGON]
    idx = 0
    for slide in prs.slides:
        for shp in iter_shapes(slide.shapes):
            if shp.has_text_frame:
                txt = shp.text_frame.text.strip()
                if txt in codes:
                    shp.text_frame.text = ""
                    try:
                        shp.auto_shape_type = shape_types[idx % len(shape_types)]
                    except:
                        pass
                    idx += 1
                    # Add label below (we add it to slide, not inside the group)
                    tb = slide.shapes.add_textbox(shp.left - Inches(0.2), shp.top + shp.height, shp.width + Inches(0.4), Inches(0.3))
                    p = tb.text_frame.paragraphs[0]
                    p.text = txt
                    p.font.size = Pt(10)
                    p.font.color.rgb = DARK_GRAY
                    tb.text_frame.word_wrap = False

    # 5. Colors & Dark theme
    for i, slide in enumerate(prs.slides):
        slide_idx = i + 1
        if slide_idx in [1, 11]:
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            rect.fill.solid()
            rect.fill.fore_color.rgb = ACCENTURE_PURPLE
            rect.line.fill.background()
            slide.shapes._spTree.insert(2, rect.element)
            
            for shp in iter_shapes(slide.shapes):
                if shp.has_text_frame:
                    for para in shp.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.color.rgb = WHITE
                            
        for shp in iter_shapes(slide.shapes):
            if getattr(shp, 'fill', None) and shp.fill.type == 1:
                try:
                    color_rgb = shp.fill.fore_color.rgb
                    if color_rgb != WHITE and color_rgb != RGBColor(255, 255, 255):
                        shp.fill.fore_color.rgb = ACCENTURE_PURPLE
                except:
                    shp.fill.fore_color.rgb = ACCENTURE_PURPLE
                    
            if shp.has_text_frame:
                txt = shp.text_frame.text
                if "Architecture Notes" in txt or "Pain Points Observed" in txt or "Business Need" in txt:
                    if getattr(shp, 'fill', None):
                        shp.fill.solid()
                        shp.fill.fore_color.rgb = WHITE
                    for p in shp.text_frame.paragraphs:
                        for r in p.runs:
                            r.font.color.rgb = DARK_GRAY

    # 9. Renumbering cascade & proper footer
    # Note: Only chip badges that are in the top-right corner are renumbered!
    for i, slide in enumerate(prs.slides):
        slide_idx = i + 1
        tb = slide.shapes.add_textbox(Inches(7.5), Inches(7.1), Inches(2.5), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = f"Karteek Tadimalla | Slide {slide_idx}"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY if slide_idx not in [1, 11] else WHITE

        for shp in iter_shapes(slide.shapes):
            if shp.has_text_frame:
                txt = shp.text_frame.text.strip()
                # Check if it's a chip badge and positioned at the top-right
                if re.match(r'^0?\d+$', txt) and len(txt) <= 2:
                    if shp.left > Inches(8) and shp.top < Inches(2):
                        # It is a large slide number badge!
                        # The original index was i (because index 1 was 02).
                        # Since we split Slide 2, index 2 became 03.
                        # So the new text should be `slide_idx:02d`.
                        p_badge = shp.text_frame.paragraphs[0]
                        p_badge.text = f"{slide_idx:02d}"
                        if p_badge.runs:
                            p_badge.runs[0].font.color.rgb = WHITE
                            p_badge.runs[0].font.bold = True

    prs.save(output_path)

if __name__ == '__main__':
    fix_presentation(
        'C:/Mortgage Calculator/ServiceNow_x_Accenture_Virtual_Agent.pptx',
        'C:/Mortgage Calculator/ServiceNow_x_Accenture_Virtual_Agent_final.pptx'
    )
