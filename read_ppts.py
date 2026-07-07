import pptx

def read_ppt(file_path, out_file):
    out_file.write(f'\n--- {file_path} ---\n')
    try:
        prs = pptx.Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            out_file.write(f'Slide {i+1}:\n')
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    out_file.write(repr(shape.text) + '\n')
    except Exception as e:
        out_file.write(f'Error: {e}\n')

with open(r'C:\Mortgage Calculator\ppt_content.txt', 'w', encoding='utf-8') as f:
    read_ppt(r'C:\Users\karte\Downloads\Virtual_Agent_Self_Service_Incident_Management_Deck.pptx', f)
    read_ppt(r'C:\Users\karte\Downloads\ServiceNow x Accenture Presentation.pptx', f)
