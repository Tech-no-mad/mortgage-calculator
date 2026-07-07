import pptx
import copy

def test():
    prs = pptx.Presentation(r'C:\Users\karte\Downloads\ServiceNow_x_Accenture_Virtual_Agent.pptx')
    slide2 = prs.slides[1]
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    for shp in slide2.shapes:
        el = shp.element
        new_el = copy.deepcopy(el)
        slide3.shapes._spTree.append(new_el)
        
    prs.save('test_dup.pptx')
    print("Success!")

test()
